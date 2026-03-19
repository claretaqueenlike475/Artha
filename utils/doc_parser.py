import os
import csv
import PyPDF2
import openpyxl
import docx
from pptx import Presentation

def parse_uploaded_file(filepath: str) -> dict:
    """
    Parse a user-uploaded file and return its content in a structure the agent can reason over.

    Dispatches to the appropriate parser based on file extension.
    Supports: .pdf, .xlsx/.xls, .csv, .docx, .pptx, and .txt.

    Args:
        filepath: Absolute path to the file in the uploads/ directory.

    Returns:
        Dict with keys: type, filename, content, char_count.
        On any parsing error, returns {"type": "error", "filename": ..., "content": str(error)}.
    """
    filename = os.path.basename(filepath)
    try:
        ext = os.path.splitext(filepath)[1].lower()
        
        # Dispatch to the appropriate parser
        if ext == ".pdf":
            content = _extract_pdf_text(filepath)
            type_ = "pdf"
        elif ext in (".xlsx", ".xls"):
            content = _extract_xlsx_tables(filepath)
            type_ = "xlsx"
        elif ext == ".csv":
            content = _extract_csv_tables(filepath)
            type_ = "csv"
        elif ext == ".docx":
            content = _extract_docx_text(filepath)
            type_ = "docx"
        elif ext == ".pptx":
            content = _extract_pptx_text(filepath)
            type_ = "pptx"
        else:
            # Fallback for .txt, .md, and other raw text files
            content = _extract_text(filepath)
            type_ = "text"

        # Compute character count for the agent's context window manager
        if isinstance(content, str):
            char_count = len(content)
        elif isinstance(content, dict):
            # For tabular data (Excel/CSV), sum the string length of all valid cells
            char_count = 0
            for sheet_name, rows in content.items():
                for row in rows:
                    for cell in row:
                        if cell is not None:
                            char_count += len(str(cell))
        else:
            char_count = 0

        return {
            "type": type_, 
            "filename": filename, 
            "content": content, 
            "char_count": char_count
        }

    except Exception as e:
        return {"type": "error", "filename": filename, "content": str(e), "char_count": 0}


def _extract_pdf_text(filepath: str) -> str:
    """Extract all text from a PDF file as a single concatenated string."""
    with open(filepath, "rb") as file_handle:
        reader = PyPDF2.PdfReader(file_handle)
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)


def _extract_xlsx_tables(filepath: str) -> dict[str, list[list]]:
    """Extract all sheets from an Excel file as a dict of 2D lists."""
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    result = {}
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows():
            row_values = [cell.value for cell in row]
            # Skip fully empty rows to reduce context noise
            if any(v is not None for v in row_values):
                rows.append(row_values)
        result[sheet_name] = rows
        
    wb.close()
    return result


def _extract_csv_tables(filepath: str) -> dict[str, list[list]]:
    """
    Extract CSV data into the same 2D list format as Excel for consistency.
    """
    rows = []
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            if any(v.strip() for v in row):  # Skip empty rows
                rows.append(row)
    # Wrap in a dict to match the multi-sheet structure of Excel
    return {"CSV_Data": rows}


def _extract_docx_text(filepath: str) -> str:
    """Extract text from Word documents."""
    doc = docx.Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])


def _extract_pptx_text(filepath: str) -> str:
    """Extract text from all slides in a PowerPoint presentation."""
    prs = Presentation(filepath)
    text_runs = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(shape.text.strip())
                
    return "\n\n".join(text_runs)


def _extract_text(filepath: str) -> str:
    """Fallback plain text reader."""
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
